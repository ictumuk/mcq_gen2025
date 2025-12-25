from django.shortcuts import render, redirect
from django.http import JsonResponse
from django.contrib.auth import login, logout
from django.contrib.auth.decorators import login_required
from django.contrib import messages
from django.views.decorators.http import require_http_methods, require_POST

from django.db.models import Max
from django.core.paginator import Paginator, EmptyPage, PageNotAnInteger
from django.utils import timezone
import json
import uuid
import PyPDF2
import docx
import pptx
from graph.g import run_mcq_generation
from .forms import RegisterForm, LoginForm, ProfileForm
from .models import Context, Question, Subject, SourceFile


def normalize_difficulty(value: str) -> str:
    """Normalize difficulty to easy/medium/hard."""
    if not value:
        return 'medium'
    value = str(value).lower()
    if value in {'easy', 'medium', 'hard'}:
        return value
    return 'medium'


def difficulty_to_bloom(difficulty: str) -> str:
    """Map normalized difficulty to Bloom level for model prompts."""
    mapping = {
        'easy': 'remember',
        'medium': 'understand',
        'hard': 'apply'
    }
    return mapping.get(normalize_difficulty(difficulty), 'understand')


def home(request):
    """Home page view"""
    return render(request, 'home.html')


@login_required
def generate_mcq(request):
    """MCQ generation page - requires login"""
    # Get 3 latest subjects for history panel (ordered by updated_at - most recently updated first)
    subjects = request.user.subjects.all().order_by('-updated_at')[:3]
    context = {
        'subjects': subjects,
    }
    return render(request, 'generate_mcq.html', context)


# ============== AUTHENTICATION VIEWS ==============

def register_view(request):
    """User registration view"""
    if request.user.is_authenticated:
        return redirect('home')
    
    if request.method == 'POST':
        form = RegisterForm(request.POST)
        if form.is_valid():
            user = form.save()
            login(request, user)
            messages.success(request, f'Chào mừng {user.username}! Tài khoản của bạn đã được tạo thành công.')
            return redirect('home')
        else:
            # Add form errors to messages
            for field, errors in form.errors.items():
                for error in errors:
                    messages.error(request, f'{error}')
    else:
        form = RegisterForm()
    
    return render(request, 'auth/register.html', {'form': form})


def login_view(request):
    """User login view"""
    if request.user.is_authenticated:
        return redirect('home')
    
    if request.method == 'POST':
        form = LoginForm(request, data=request.POST)
        if form.is_valid():
            user = form.get_user()
            login(request, user)
            
            # Handle remember me
            if not form.cleaned_data.get('remember_me'):
                request.session.set_expiry(0)  # Session expires when browser closes
            
            messages.success(request, f'Chào mừng trở lại, {user.username}!')
            
            # Redirect to next page if specified
            next_url = request.GET.get('next', 'home')
            return redirect(next_url)
        else:
            messages.error(request, 'Tên đăng nhập hoặc mật khẩu không đúng.')
    else:
        form = LoginForm()
    
    return render(request, 'auth/login.html', {'form': form})


@require_http_methods(["GET", "POST"])
def logout_view(request):
    """User logout view"""
    if request.user.is_authenticated:
        logout(request)
        messages.info(request, 'Bạn đã đăng xuất thành công.')
    return redirect('home')


# ============== PROFILE VIEWS ==============

@login_required
def profile_view(request):
    """User profile view"""
    if request.method == 'POST':
        form = ProfileForm(request.POST, instance=request.user)
        if form.is_valid():
            form.save()
            messages.success(request, 'Thông tin cá nhân đã được cập nhật.')
            return redirect('profile')
    else:
        form = ProfileForm(instance=request.user)
    
    return render(request, 'auth/profile.html', {'form': form})


@login_required
def dashboard_view(request):
    """User dashboard with history and pagination"""
    # Get all user's subjects ordered by updated_at (most recently updated first)
    all_subjects = request.user.subjects.all().order_by('-updated_at')
    
    # Paginate: 10 items per page
    paginator = Paginator(all_subjects, 10)
    page = request.GET.get('page', 1)
    
    try:
        subjects = paginator.page(page)
    except PageNotAnInteger:
        subjects = paginator.page(1)
    except EmptyPage:
        subjects = paginator.page(paginator.num_pages)
    
    # Calculate total questions from all subjects (not just current page)
    total_questions = sum(subj.questions.count() for subj in all_subjects)
    
    context = {
        'subjects': subjects,
        'total_questions': total_questions,
        'credits': request.user.credits,
        'paginator': paginator,
    }
    return render(request, 'auth/dashboard.html', context)

# ============== API VIEWS FOR QUESTION MANAGEMENT ==============

@login_required
@require_POST
def api_create_question(request):
    """
    API endpoint to create a new question manually.
    Expects JSON body with question data.
    """
    try:
        data = json.loads(request.body)
        
        # Validate required fields
        content = data.get('content', '').strip()
        options = data.get('options', [])
        correct_answer = data.get('correct_answer', 'A')
        difficulty = normalize_difficulty(data.get('difficulty', 'medium'))
        bloom_level = data.get('bloom_level') or difficulty_to_bloom(difficulty)
        subject_name = data.get('subject', '').strip()
        
        if not content:
            return JsonResponse({
                'success': False, 
                'error': 'Nội dung câu hỏi không được để trống'
            }, status=400)
        
        if len(options) < 2:
            return JsonResponse({
                'success': False, 
                'error': 'Cần ít nhất 2 đáp án'
            }, status=400)
        
        # Get or create subject based on subject_id or subject_name
        subject_id = data.get('subject_id')
        # Normalize subject_id: convert empty string to None
        if subject_id == '' or subject_id is None:
            subject_id = None
        subject_obj = None
        
        if subject_id:
            # Use existing subject if subject_id is provided
            # If subject_id is provided, we must find it - don't fallback to name search
            try:
                subject_obj = Subject.objects.get(id=subject_id, user=request.user)
            except Subject.DoesNotExist:
                return JsonResponse({
                    'success': False, 
                    'error': f'Không tìm thấy môn học với ID: {subject_id}'
                }, status=404)
            except (ValueError, TypeError):
                # Invalid UUID format
                return JsonResponse({
                    'success': False, 
                    'error': 'ID môn học không hợp lệ'
                }, status=400)
        
        if not subject_obj and subject_name:
            # Find existing subject with matching subject name, or create new one
            # Only search by name if subject_id was not provided
            subject_obj = Subject.objects.filter(
                user=request.user,
                subject=subject_name
            ).first()
            
            if not subject_obj:
                # Create new subject with the specified name
                subject_obj = Subject.objects.create(
                    user=request.user,
                    title=f'Gen MCQ - {subject_name}',
                    source_type='text',
                    subject=subject_name,
                    status='completed',
                    source_text='Manual questions',
                    difficulty=difficulty,
                    bloom_level=bloom_level.capitalize() if bloom_level else 'Understand'
                )
        elif not subject_obj:
            # Fallback to default subject if no subject_id and no subject_name provided
            subject_obj, created = Subject.objects.get_or_create(
                user=request.user,
                title='Câu hỏi thủ công',
                source_type='text',
                defaults={
                    'status': 'completed',
                    'source_text': 'Manual questions',
                    'difficulty': difficulty,
                    'bloom_level': bloom_level.capitalize() if bloom_level else 'Understand'
                }
            )
        
        # Determine the next order number
        max_order = subject_obj.questions.aggregate(
            max_order=Max('order')
        )['max_order'] or 0
        
        # Create the question
        question = Question.objects.create(
            subject=subject_obj,
            question_type='mcq',
            stem=content,
            difficulty=difficulty,
            options=options,
            correct_answer=correct_answer,
            explanation=data.get('explanation', ''),
            reasoning={
                'bloom_level': bloom_level,
                'difficulty': difficulty
            },
            user_edited=True,
            order=max_order + 1
        )
        
        # Update subject's updated_at to reflect the latest question creation time
        subject_obj.updated_at = timezone.now()
        subject_obj.save(update_fields=['updated_at'])
        
        # Return the created question data
        return JsonResponse({
            'success': True,
            'message': 'Câu hỏi đã được tạo thành công',
            'question': {
                'id': str(question.id),
                'stem': question.stem,
                'options': question.options,
                'correct_answer': question.correct_answer,
                'explanation': question.explanation,
                'bloom_level': question.reasoning.get('bloom_level', bloom_level),
                'difficulty': question.difficulty,
                'order': question.order,
                'user_edited': question.user_edited
            },
            'user_credits': request.user.credits
        })
        
    except json.JSONDecodeError:
        return JsonResponse({
            'success': False, 
            'error': 'Dữ liệu JSON không hợp lệ'
        }, status=400)
    except Exception as e:
        return JsonResponse({
            'success': False, 
            'error': str(e)
        }, status=500)


@login_required
@require_POST
def api_update_question(request, question_id):
    """
    API endpoint to update an existing question.
    """
    try:
        data = json.loads(request.body)
        
        # Get the question (ensure user owns it)
        question = Question.objects.select_related('subject').get(
            id=question_id,
            subject__user=request.user
        )
        
        # Store original data before editing
        if not question.original_data:
            question.original_data = {
                'stem': question.stem,
                'options': question.options,
                'correct_answer': question.correct_answer,
                'explanation': question.explanation
            }
        
        # Update fields
        if 'content' in data:
            question.stem = data['content']
        if 'options' in data:
            question.options = data['options']
        if 'correct_answer' in data:
            question.correct_answer = data['correct_answer']
        if 'explanation' in data:
            question.explanation = data['explanation']
        if 'bloom_level' in data or 'difficulty' in data:
            new_difficulty = normalize_difficulty(data.get('difficulty', question.difficulty))
            new_bloom = data.get('bloom_level', question.reasoning.get('bloom_level')) or difficulty_to_bloom(new_difficulty)
            question.difficulty = new_difficulty
            question.reasoning = {
                **question.reasoning,
                'bloom_level': new_bloom,
                'difficulty': new_difficulty
            }
        
        question.user_edited = True
        question.save()
        
        return JsonResponse({
            'success': True,
            'message': 'Câu hỏi đã được cập nhật',
            'question': {
                'id': str(question.id),
                'stem': question.stem,
                'options': question.options,
                'correct_answer': question.correct_answer,
                'explanation': question.explanation,
                'user_edited': question.user_edited
            }
        })
        
    except Question.DoesNotExist:
        return JsonResponse({
            'success': False, 
            'error': 'Không tìm thấy câu hỏi'
        }, status=404)
    except Exception as e:
        return JsonResponse({
            'success': False, 
            'error': str(e)
        }, status=500)


@login_required
@require_POST
def api_delete_question(request, question_id):
    """
    API endpoint to delete a question.
    """
    try:
        question = Question.objects.select_related('subject').get(
            id=question_id,
            subject__user=request.user
        )
        
        question.delete()
        
        return JsonResponse({
            'success': True,
            'message': 'Câu hỏi đã được xóa'
        })
        
    except Question.DoesNotExist:
        return JsonResponse({
            'success': False, 
            'error': 'Không tìm thấy câu hỏi'
        }, status=404)
    except Exception as e:
        return JsonResponse({
            'success': False, 
            'error': str(e)
        }, status=500)


@login_required
def api_get_subjects_list(request):
    """
    API endpoint to get list of subjects for the current user.
    Returns both subject_id and subject name for each subject.
    """
    try:
        # Get subjects with their IDs and names, grouped by subject name
        # For each unique subject name, return the most recent subject's ID
        subjects_data = []
        seen_names = set()
        
        # Get all subjects ordered by created_at (newest first)
        all_subjects = Subject.objects.filter(
            user=request.user,
            subject__isnull=False
        ).exclude(
            subject=''
        ).order_by('-created_at')
        
        for subj in all_subjects:
            if subj.subject not in seen_names:
                seen_names.add(subj.subject)
                subjects_data.append({
                    'id': str(subj.id),
                    'name': subj.subject
                })
        
        # Sort by name for display
        subjects_data.sort(key=lambda x: x['name'])
        
        return JsonResponse({
            'success': True,
            'subjects': subjects_data,
            'count': len(subjects_data)
        })
        
    except Exception as e:
        return JsonResponse({
            'success': False, 
            'error': str(e)
        }, status=500)


@login_required
def api_get_subjects_history(request):
    """
    API endpoint to get subjects list for history panel (with full info).
    Returns latest 3 subjects with question counts, ordered by updated_at.
    """
    try:
        # Order by updated_at to show most recently updated subjects first
        subjects = Subject.objects.filter(
            user=request.user
        ).order_by('-updated_at')[:3]
        
        subjects_data = []
        for subj in subjects:
            subjects_data.append({
                'id': str(subj.id),
                'title': subj.title or 'Không có tiêu đề',
                'subject': subj.subject or '',
                'difficulty': subj.difficulty or 'medium',
                'question_count': subj.questions.count(),
                'created_at': subj.created_at.isoformat(),
                'updated_at': subj.updated_at.isoformat()
            })
        
        return JsonResponse({
            'success': True,
            'subjects': subjects_data,
            'count': len(subjects_data)
        })
        
    except Exception as e:
        return JsonResponse({
            'success': False, 
            'error': str(e)
        }, status=500)


@login_required
def api_get_questions(request):
    """
    API endpoint to get all questions for the current user.
    Questions are ordered by created_at (newest first) to show most recently created questions first.
    """
    try:
        subject_id = request.GET.get('subject_id')
        
        if subject_id:
            questions = Question.objects.filter(
                subject_id=subject_id,
                subject__user=request.user
            ).order_by('-created_at')  # Newest questions first
        else:
            # Nếu không truyền subject_id: lấy subject mới nhất được cập nhật của user
            latest_subj = Subject.objects.filter(user=request.user).order_by('-updated_at').first()
            if not latest_subj:
                questions = Question.objects.none()
            else:
                questions = Question.objects.filter(
                    subject=latest_subj
                ).order_by('-created_at')  # Newest questions first
        
        questions_data = []
        for q in questions:
            difficulty = normalize_difficulty(getattr(q, 'difficulty', 'medium'))
            bloom = q.reasoning.get('bloom_level') if isinstance(q.reasoning, dict) else None
            bloom = bloom or difficulty_to_bloom(difficulty)
            questions_data.append({
                'id': str(q.id),
                'stem': q.stem,
                'options': q.options,
                'correct_answer': q.correct_answer,
                'explanation': q.explanation,
                'bloom_level': bloom,
                'difficulty': difficulty,
                'order': q.order,
                'user_edited': q.user_edited,
                'created_at': q.created_at.isoformat()
            })
        
        return JsonResponse({
            'success': True,
            'questions': questions_data,
            'count': len(questions_data)
        })
        
    except Exception as e:
        return JsonResponse({
            'success': False, 
            'error': str(e)
        }, status=500)


@login_required
@require_POST  
def api_clear_all_questions(request):
    """
    API endpoint to clear all manual questions for the current user.
    """
    try:
        # Delete all questions from manual subject
        deleted_count, _ = Question.objects.filter(
            subject__user=request.user,
            subject__title='Câu hỏi thủ công'
        ).delete()
        
        return JsonResponse({
            'success': True,
            'message': f'Đã xóa {deleted_count} câu hỏi',
            'deleted_count': deleted_count
        })
        
    except Exception as e:
        return JsonResponse({
            'success': False, 
            'error': str(e)
        }, status=500)

@login_required
@require_POST
def api_delete_subject(request, subject_id):
    """
    API endpoint to delete a subject and all its questions.
    """
    try:
        subject = Subject.objects.get(
            id=subject_id,
            user=request.user
        )
        
        title = subject.title
        subject.delete()
        
        return JsonResponse({
            'success': True,
            'message': f'Đã xóa môn học "{title}"'
        })
        
    except Subject.DoesNotExist:
        return JsonResponse({
            'success': False,
            'error': 'Không tìm thấy môn học'
        }, status=404)
    except Exception as e:
        return JsonResponse({
            'success': False,
            'error': str(e)
        }, status=500)

def extract_text_from_sourcefile(sf: SourceFile) -> str:
    """
    Extract text from SourceFile based on file_type.
    Supports: pdf (PyPDF2), docx (python-docx), pptx (python-pptx), txt.
    """
    if not sf.file:
        raise ValueError("SourceFile không có file được đính kèm")
    
    path = sf.file.path
    ext = sf.file_type

    try:
        if ext == 'pdf':
            text_parts = []
            with open(path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    txt = page.extract_text() or ''
                    text_parts.append(txt)
            result = '\n'.join(text_parts)
            if not result.strip():
                raise ValueError("Không thể trích xuất văn bản từ file PDF")
            return result

        if ext == 'docx':
            doc = docx.Document(path)
            result = '\n'.join(p.text for p in doc.paragraphs)
            if not result.strip():
                raise ValueError("Không thể trích xuất văn bản từ file DOCX")
            return result

        if ext == 'pptx':
            prs = pptx.Presentation(path)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)
            result = '\n'.join(texts)
            if not result.strip():
                raise ValueError("Không thể trích xuất văn bản từ file PPTX")
            return result

        if ext == 'txt':
            # Try multiple encodings
            encodings = ['utf-8', 'utf-16', 'latin-1', 'cp1252']
            for encoding in encodings:
                try:
                    with open(path, 'r', encoding=encoding, errors='ignore') as f:
                        result = f.read()
                        if result.strip():
                            return result
                except (UnicodeDecodeError, UnicodeError):
                    continue
            # Fallback: read as binary and decode with errors='ignore'
            with open(path, 'rb') as f:
                return f.read().decode('utf-8', errors='ignore')

        raise ValueError(f"Định dạng file '{ext}' chưa được hỗ trợ để trích xuất.")
    
    except FileNotFoundError:
        raise ValueError(f"Không tìm thấy file tại đường dẫn: {path}")
    except PermissionError:
        raise ValueError(f"Không có quyền truy cập file: {path}")
    except Exception as e:
        raise ValueError(f"Lỗi khi trích xuất văn bản từ file: {str(e)}")

@login_required
def api_generate_mcq(request):
    """
    API endpoint to generate a new MCQ.
    """
    try:
        data = json.loads(request.body)
        source_type = data.get('source_type', 'text')
        text = data.get('text', '')
        subject = data.get('subject', '')
        topic = data.get('topic', '')
        difficulty = normalize_difficulty(data.get('level', data.get('difficulty', 'medium')))
        bloom_level = difficulty_to_bloom(difficulty)
        # Đồng bộ key: frontend gửi number_contexts, fallback number_questions
        number_contexts = data.get('number_contexts')
        if number_contexts is None:
            number_contexts = data.get('number_questions', 3)
        key_point = data.get('key_point', '')
        exercises = data.get('exercises', '')
        model = "gemini-2.5-flash"
        max_iterations = 2
        max_workers = 1
        delay_seconds = 5.0
        sf = None

        # If file source, try to extract text from file if not already stored
        if source_type == 'file':
            source_file_id = data.get('source_file_id')
            if not source_file_id:
                return JsonResponse({'success': False, 'error': 'Thiếu source_file_id'}, status=400)
            try:
                sf = SourceFile.objects.get(id=source_file_id, user=request.user)
            except SourceFile.DoesNotExist:
                return JsonResponse({'success': False, 'error': 'Không tìm thấy file nguồn'}, status=404)
            if sf.extracted_text:
                text = sf.extracted_text
            else:
                try:
                    text_extracted = extract_text_from_sourcefile(sf)
                except ImportError as ie:
                    return JsonResponse({'success': False, 'error': f'Thiếu thư viện để đọc file: {ie}'}, status=500)
                except Exception as e:
                    return JsonResponse({'success': False, 'error': f'Lỗi xử lý file: {e}'}, status=500)
                if not text_extracted.strip():
                    return JsonResponse({'success': False, 'error': 'Không trích xuất được nội dung từ file'}, status=400)
                sf.extracted_text = text_extracted
                sf.save(update_fields=['extracted_text'])
                text = text_extracted

        if not text.strip():
            return JsonResponse({'success': False, 'error': 'Thiếu nội dung văn bản để tạo câu hỏi'}, status=400)

        # Run generation with thread_id
        thread_id = str(uuid.uuid4())
        result = run_mcq_generation(
            text=text,
            subject=subject,
            topic=topic,
            bloom_level=bloom_level,
            number_contexts=number_contexts,
            key_point=key_point,
            exercises=exercises,
            model=model,
            max_iterations=max_iterations,
            max_workers=max_workers,
            delay_seconds=delay_seconds
        )

        contexts_result = result.get('contexts', [])
        mcqs = result.get('mcqs', [])

        # Get or create Subject
        subject_id = data.get('subject_id')
        if subject_id:
            # Use existing subject if subject_id is provided
            try:
                subject_obj = Subject.objects.get(id=subject_id, user=request.user)
                # Update subject with new generation info
                subject_obj.number_questions = subject_obj.questions.count() + len(mcqs)
                subject_obj.credits_used = (subject_obj.credits_used or 0) + len(mcqs)
                subject_obj.save(update_fields=['number_questions', 'credits_used'])
            except Subject.DoesNotExist:
                # Subject not found, create new one
                subject_obj = None
        else:
            subject_obj = None
        
        if not subject_obj:
            # Create new Subject
            subject_obj = Subject.objects.create(
                user=request.user,
                title=f"Gen MCQ - {subject or 'Untitled'}",
                source_type=source_type,
                source_file=sf if source_type == 'file' else None,
                source_text=text if source_type == 'text' else '',
                subject=subject,
                topic=topic,
                difficulty=difficulty,
                bloom_level=bloom_level.capitalize() if bloom_level else 'Understand',
                number_questions=len(mcqs),
                key_points=key_point if hasattr(Subject, 'key_points') else '',
                exercises=exercises,
                status='completed',
                thread_id=thread_id,
                current_stage=result.get('current_stage', 'complete'),
                iteration_count=result.get('mcq_iteration', 0),
                config={
                    "model": model,
                    "max_iterations": max_iterations,
                    "max_workers": max_workers,
                    "delay_seconds": delay_seconds,
                    "difficulty": difficulty,
                    "bloom_level": bloom_level
                },
                credits_used=len(mcqs)
            )

        # Save contexts to DB
        context_objs = []
        for idx, ctx_item in enumerate(contexts_result):
            ctx_data = ctx_item if isinstance(ctx_item, dict) else getattr(ctx_item, '__dict__', {}) or {}
            context_objs.append(
                Context.objects.create(
                    subject=subject_obj,
                    content=ctx_data.get('context', ''),
                    original_content=ctx_data.get('original_content', ctx_data.get('context', '')),
                    review_feedback=ctx_data.get('review', ''),
                    suggestions=ctx_data.get('suggestions', []),
                    is_approved=ctx_data.get('is_approved', False),
                    iteration_count=ctx_data.get('iteration_count', 0),
                    order=idx
                )
            )
        context_map = {idx: ctx for idx, ctx in enumerate(context_objs)}

        # Save questions
        questions_payload = []
        for idx, mcq_item in enumerate(mcqs):
            mcq_obj = mcq_item.get('mcq') if isinstance(mcq_item, dict) else getattr(mcq_item, 'mcq', mcq_item)
            if isinstance(mcq_obj, dict):
                q_data = mcq_obj.get('question', mcq_obj)
            else:
                q_data = getattr(mcq_obj, 'question', mcq_obj)

            # Normalize question payload
            if hasattr(q_data, 'model_dump'):
                q_data = q_data.model_dump()
            elif hasattr(q_data, 'dict'):
                q_data = q_data.dict()
            elif hasattr(q_data, '__dict__'):
                q_data = {k: v for k, v in q_data.__dict__.items() if not k.startswith('_')}
            elif not isinstance(q_data, dict):
                q_data = {}

            stem = q_data.get('stem', '')
            options_data = q_data.get('options', {})
            if isinstance(options_data, dict):
                options_list = options_data.get('options', [])
            else:
                options_list = options_data or []

            correct_answer = q_data.get('correct_answer', '')
            normalized_options = []
            for opt in options_list:
                if isinstance(opt, dict):
                    opt_id = opt.get('id')
                    opt_text = opt.get('text', '')
                    is_correct_flag = opt.get('is_correct')
                else:
                    opt_id = getattr(opt, 'id', None)
                    opt_text = getattr(opt, 'text', '')
                    is_correct_flag = getattr(opt, 'is_correct', None)
                normalized_options.append({
                    "id": opt_id,
                    "text": opt_text,
                    "is_correct": is_correct_flag if is_correct_flag is not None else str(opt_id).lower() == str(correct_answer).lower()
                })

            reasoning = q_data.get('reasoning', {}) or {}
            if not isinstance(reasoning, dict):
                reasoning = {}
            reasoning = {
                **reasoning,
                "bloom_level": bloom_level,
                "difficulty": difficulty
            }
            explanation = reasoning.get('answer_justification', q_data.get('explanation', ''))

            review_feedback = ''
            suggestions = []
            is_approved = False
            context_ref = None
            if isinstance(mcq_item, dict):
                review_feedback = mcq_item.get('review', '')
                suggestions = mcq_item.get('suggestions', [])
                is_approved = mcq_item.get('is_approved', False)
                context_ref = context_map.get(mcq_item.get('context_index'))
            elif hasattr(mcq_item, 'context_index'):
                context_ref = context_map.get(getattr(mcq_item, 'context_index'))

            q_instance = Question.objects.create(
                subject=subject_obj,
                context=context_ref,
                question_type='mcq',
                stem=stem,
                difficulty=difficulty,
                options=normalized_options,
                correct_answer=correct_answer,
                explanation=explanation,
                reasoning=reasoning,
                review_feedback=review_feedback,
                suggestions=suggestions,
                is_approved=is_approved,
                user_edited=False,
                order=idx
            )

            questions_payload.append({
                'id': str(q_instance.id),
                'stem': q_instance.stem,
                'options': q_instance.options,
                'correct_answer': q_instance.correct_answer,
                'explanation': q_instance.explanation,
                'bloom_level': q_instance.reasoning.get('bloom_level', bloom_level),
                'difficulty': q_instance.difficulty,
                'order': q_instance.order,
                'user_edited': q_instance.user_edited
            })

        contexts_payload = [{
            'id': str(ctx.id),
            'content': ctx.content,
            'order': ctx.order,
            'is_approved': ctx.is_approved,
            'review_feedback': ctx.review_feedback,
            'suggestions': ctx.suggestions
        } for ctx in context_objs]

        # Update subject's updated_at to reflect the latest question creation time
        if len(mcqs) > 0:
            subject_obj.updated_at = timezone.now()
            subject_obj.save(update_fields=['updated_at'])
        
        # Deduct credits from user after successful question generation
        # Deduct 1 credit per generation request (not per question)
        if len(mcqs) > 0:
            # Refresh user from database to get latest credits
            request.user.refresh_from_db()
            # Deduct 1 credit per generation request
            credits_deducted = request.user.use_credits(1)
            if not credits_deducted:
                # If not enough credits, still return success but log warning
                # (Questions are already created, so we don't rollback)
                pass
            # Refresh again to get updated credits value
            request.user.refresh_from_db()

        return JsonResponse({
            'success': True,
            'message': 'MCQ đã được tạo thành công',
            'subject_id': str(subject_obj.id),
            'thread_id': thread_id,
            'difficulty': difficulty,
            'contexts': contexts_payload,
            'questions': questions_payload,
            'user_credits': request.user.credits
        })
    except Exception as e:
        return JsonResponse({
            'success': False, 
            'error': str(e)
        }, status=500)


@login_required
@require_POST
def api_upload_source(request):
    """
    Upload source file (pdf, docx, pptx, txt). Note: chưa trích xuất nội dung.
    """
    try:
        if 'file' not in request.FILES:
            return JsonResponse({'success': False, 'error': 'Không có file'}, status=400)

        file_obj = request.FILES['file']
        filename = file_obj.name
        ext = filename.split('.')[-1].lower()
        allowed = {'pdf': 'pdf', 'docx': 'docx', 'pptx': 'pptx', 'txt': 'txt'}
        if ext not in allowed:
            return JsonResponse({'success': False, 'error': 'Định dạng không hỗ trợ'}, status=400)

        sf = SourceFile.objects.create(
            user=request.user,
            file_name=filename,
            file_type=allowed[ext],
            file=file_obj,
            extracted_text=''  
        )

        return JsonResponse({
            'success': True,
            'source_file_id': str(sf.id),
            'file_name': sf.file_name,
            'file_type': sf.file_type
        })
    except Exception as e:
        return JsonResponse({'success': False, 'error': str(e)}, status=500)